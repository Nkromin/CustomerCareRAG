"""Generate a project overview presentation (.pptx) for CustomerCareRAG."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def _add_title_slide(prs, title, subtitle=""):
    """Add a title slide to the presentation."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle
    return slide


def _add_content_slide(prs, title, bullet_groups):
    """Add a content slide with grouped bullet points.

    Parameters
    ----------
    bullet_groups : list[tuple[str, list[str]]]
        Each tuple is (section_heading, [bullet_items]).
        If section_heading is empty the bullets are top-level.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    tf = slide.placeholders[1].text_frame
    tf.clear()

    first = True
    for heading, items in bullet_groups:
        if heading:
            if not first:
                p = tf.add_paragraph()
                p.space_before = Pt(6)
            else:
                p = tf.paragraphs[0]
                first = False
            p.text = heading
            p.font.bold = True
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
            p.alignment = PP_ALIGN.LEFT

        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 1
            p.font.size = Pt(16)
            p.alignment = PP_ALIGN.LEFT
            if first:
                first = False

    return slide


def generate():
    """Create the presentation and save it to docs/."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 11 — Local Testing & Demo Scenarios ---
    demo_queries = [
        '"What is sick leave policy?"',
        '"Can I take leave for sinus infection?"',
        '"Create a ticket for laptop issue"',
        '"Check my leave balance"',
    ]
    validations = [
        "Correct routing (retrieve / tool / general)",
        "Policy-grounded responses",
        "Tool execution outputs",
    ]
    _add_content_slide(
        prs,
        "Local Testing & Demo Scenarios",
        [
            ("Run with Streamlit UI and ask:", demo_queries),
            ("Validate:", validations),
        ],
    )

    # --- Slide 12 — Key Benefits & Future Work ---
    benefits = [
        "Faster HR query handling",
        "Better employee self-service",
        "Consistent policy interpretation",
        "Extensible tool-based operations",
    ]
    future = [
        "Real HRMS API integration",
        "Authentication + role-based answers",
        "Observability dashboards",
        "Source citations per response",
        "Production deployment and monitoring",
    ]
    _add_content_slide(
        prs,
        "Key Benefits & Future Work",
        [
            ("Benefits", benefits),
            ("Future Enhancements", future),
        ],
    )

    output_path = "docs/CustomerCareRAG_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")


if __name__ == "__main__":
    generate()
